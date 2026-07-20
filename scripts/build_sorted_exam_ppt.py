import argparse
import json
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

DEFAULT_SLIDE_W = 7.5
DEFAULT_SLIDE_H = 10.833333


def parse_ranges(text):
    ranges = []
    if not text:
        return ranges
    for part in text.split(','):
        part = part.strip()
        if not part:
            continue
        if '-' in part:
            a, b = part.split('-', 1)
            ranges.append((int(a), int(b)))
        else:
            n = int(part)
            ranges.append((n, n))
    return ranges


def in_ranges(value, ranges):
    return any(a <= value <= b for a, b in ranges)


def resolve_image_path(raw_path, cwd, archive_root=None):
    path = Path(raw_path)
    if path.exists():
        return path
    candidates = []
    if not path.is_absolute():
        candidates.append(cwd / path)
        if archive_root:
            candidates.append(archive_root / path)
    if archive_root:
        try:
            candidates.append(archive_root / path.relative_to(cwd))
        except ValueError:
            pass
        marker = cwd.name + '\\'
        text = str(path)
        idx = text.lower().find(marker.lower())
        if idx >= 0:
            candidates.append(archive_root / text[idx + len(marker):])
    for candidate in candidates:
        if candidate.exists():
            return candidate
    raise FileNotFoundError(f'Image not found: {raw_path}')


def fit_to_box(iw, ih, max_w, max_h):
    scale = min(max_w / iw, max_h / ih)
    return int(iw * scale), int(ih * scale)


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_white_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)


def add_chapter_divider(prs, chapter, slide_w, slide_h):
    slide = blank_slide(prs)
    set_white_bg(slide)
    band = slide.shapes.add_shape(1, 0, Inches(4.53), slide_w, Inches(1.75))
    band.fill.solid()
    band.fill.fore_color.rgb = RGBColor(220, 30, 42)
    band.line.color.rgb = RGBColor(220, 30, 42)
    box = slide.shapes.add_textbox(Inches(0.45), Inches(4.92), slide_w - Inches(0.9), Inches(0.75))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = f'Chapter {chapter}'
    r.font.name = 'Arial'
    r.font.size = Pt(40)
    r.font.bold = True
    r.font.color.rgb = RGBColor(255, 255, 255)


def add_label(slide, label, x, y, w, h, font_size):
    box = slide.shapes.add_textbox(x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0, 0, 0)
    box.line.color.rgb = RGBColor(0, 0, 0)
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label or ''
    r.font.name = 'Arial'
    r.font.size = Pt(font_size)
    r.font.bold = True
    r.font.color.rgb = RGBColor(255, 255, 255)


def rects_overlap(a, b):
    ax, ay, aw, ah = a
    bx, by, bw, bh = b
    return ax < bx + bw and ax + aw > bx and ay < by + bh and ay + ah > by


def build_deck(args):
    cwd = Path.cwd()
    archive_root = Path(args.archive_root).resolve() if args.archive_root else None
    input_path = Path(args.input_json)
    items = json.loads(input_path.read_text(encoding='utf-8-sig'))
    if args.source_id:
        items = [item for item in items if item.get('source_id') == args.source_id]

    ranges = parse_ranges(args.mcq_ranges)
    for item in items:
        qno = int(item['question_no_in_file'])
        item['_type'] = 'mcq' if in_ranges(qno, ranges) else 'written'

    slide_w = Inches(args.slide_width)
    slide_h = Inches(args.slide_height)
    top = Inches(args.top)
    bottom = Inches(args.bottom)
    gap = Inches(args.gap)
    label_band = Inches(args.label_band_height)
    label_w = Inches(args.label_width)
    label_h = Inches(args.label_height)
    label_x = slide_w - label_w - Inches(args.label_right_margin)
    image_w = slide_w

    def image_size(item):
        path = resolve_image_path(item['image_path'], cwd, archive_root)
        with Image.open(path) as image:
            iw, ih = image.size
        w, h = fit_to_box(iw, ih, image_w, slide_h - top - bottom - label_band)
        return path, w, h

    prs = Presentation()
    prs.slide_width = slide_w
    prs.slide_height = slide_h

    chapters = sorted({int(item['chapter']) for item in items})
    for chapter in chapters:
        chapter_items = [item for item in items if int(item['chapter']) == chapter]
        if not chapter_items:
            continue
        add_chapter_divider(prs, chapter, slide_w, slide_h)
        mcq = sorted([item for item in chapter_items if item['_type'] == 'mcq'], key=lambda i: int(i['question_no_in_file']))
        written = sorted([item for item in chapter_items if item['_type'] == 'written'], key=lambda i: int(i['question_no_in_file']))

        slide = None
        y = top
        for item in mcq + written:
            path, w, h = image_size(item)
            needed = label_band + h
            if slide is None or y + needed > slide_h - bottom:
                slide = blank_slide(prs)
                set_white_bg(slide)
                y = top
            x = max(0, (image_w - w) // 2)
            slide.shapes.add_picture(str(path), x, y + label_band, width=w, height=h)
            add_label(slide, item.get('source_label', ''), label_x, y, label_w, label_h, args.label_font_size)
            y += needed + gap

    prs.save(args.out)
    return args.out


def validate_deck(pptx_path):
    prs = Presentation(pptx_path)
    labels = set()
    # Source labels vary by project; treat black text boxes with text as labels unless they are Chapter dividers.
    overlap_count = 0
    picture_oob = 0
    label_oob = 0
    sw, sh = prs.slide_width, prs.slide_height
    for slide in prs.slides:
        pics = [(shape.left, shape.top, shape.width, shape.height) for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE]
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                if shape.left < 0 or shape.top < 0 or shape.left + shape.width > sw or shape.top + shape.height > sh:
                    picture_oob += 1
            if getattr(shape, 'has_text_frame', False) and shape.text.strip() and not shape.text.strip().startswith('Chapter '):
                labels.add(shape.text.strip())
                rect = (shape.left, shape.top, shape.width, shape.height)
                if shape.left < 0 or shape.top < 0 or shape.left + shape.width > sw or shape.top + shape.height > sh:
                    label_oob += 1
                if any(rects_overlap(rect, pic) for pic in pics):
                    overlap_count += 1
    return {
        'slides': len(prs.slides),
        'source_labels': sorted(labels),
        'label_picture_overlaps': overlap_count,
        'picture_out_of_bounds': picture_oob,
        'label_out_of_bounds': label_oob,
    }


def main():
    parser = argparse.ArgumentParser(description='Build a chapter-sorted exam PPT from classified question JSON.')
    parser.add_argument('--input-json', required=True)
    parser.add_argument('--out', required=True)
    parser.add_argument('--source-id')
    parser.add_argument('--archive-root')
    parser.add_argument('--mcq-ranges', default='')
    parser.add_argument('--slide-width', type=float, default=DEFAULT_SLIDE_W)
    parser.add_argument('--slide-height', type=float, default=DEFAULT_SLIDE_H)
    parser.add_argument('--top', type=float, default=0.10)
    parser.add_argument('--bottom', type=float, default=0.12)
    parser.add_argument('--gap', type=float, default=0.08)
    parser.add_argument('--label-band-height', type=float, default=0.40)
    parser.add_argument('--label-width', type=float, default=1.45)
    parser.add_argument('--label-height', type=float, default=0.34)
    parser.add_argument('--label-right-margin', type=float, default=0.21)
    parser.add_argument('--label-font-size', type=float, default=14)
    parser.add_argument('--no-validate', action='store_true')
    args = parser.parse_args()

    out = build_deck(args)
    print(out)
    if not args.no_validate:
        print(json.dumps(validate_deck(out), ensure_ascii=False, indent=2))


if __name__ == '__main__':
    main()
